"""
Version alternative de format2_methods avec fallback automatique
Essaie d'abord l'envoi direct, puis fallback sur extraction si échec
"""
import os
import json
import io
from google import genai
from google.genai import types
from typing import Union
from fastapi import UploadFile
from schemas.format2_schema import Format2SFD
from datetime import datetime

# Configuration Gemini
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
GEMINI_MODEL = 'gemini-2.5-flash-lite'
_client = genai.Client(api_key=GOOGLE_API_KEY)


async def extract_from_file(file: UploadFile) -> dict:
    """
    Extraire le contenu d'un fichier avec stratégie de fallback automatique

    Stratégie :
    1. Essayer d'envoyer directement à Gemini avec le MIME type
    2. Si échec (Unsupported MIME), extraire le texte puis réessayer

    Args:
        file: Fichier uploadé

    Returns:
        dict: Contenu extrait
    """
    file_content = await file.read()
    file_extension = file.filename.split('.')[-1].lower()

    prompt = """
    Analyse ce document et extrais toutes les informations relatives à une spécification fonctionnelle détaillée (SFD).
    Identifie et extrais:
    - Le contexte du projet et ses objectifs
    - La description générale du système
    - Les modules fonctionnels et leurs fonctions détaillées
    - Les exigences non fonctionnelles (performance, sécurité, disponibilité, scalabilité)
    - Les contraintes techniques
    - Tous les détails pertinents pour un SFD

    Retourne un texte structuré avec toutes ces informations.
    """

    # STRATÉGIE 1 : Essayer l'envoi direct
    try:
        if file_extension in ['jpg', 'jpeg', 'png', 'gif', 'webp']:
            mime_type = f"image/{file_extension}"
            response = _client.models.generate_content(
                model=GEMINI_MODEL,
                contents=[prompt, types.Part.from_bytes(data=file_content, mime_type=mime_type)]
            )
            return {"extracted_text": response.text}

        elif file_extension == 'pdf':
            response = _client.models.generate_content(
                model=GEMINI_MODEL,
                contents=[prompt, types.Part.from_bytes(data=file_content, mime_type="application/pdf")]
            )
            return {"extracted_text": response.text}

        elif file_extension in ['txt', 'md']:
            text_content = file_content.decode('utf-8')
            response = _client.models.generate_content(
                model=GEMINI_MODEL,
                contents=[prompt, text_content]
            )
            return {"extracted_text": response.text}

        elif file_extension in ['docx', 'doc']:
            # Essayer d'abord avec le MIME type Word
            mime_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            try:
                response = _client.models.generate_content(
                    model=GEMINI_MODEL,
                    contents=[prompt, types.Part.from_bytes(data=file_content, mime_type=mime_type)]
                )
                print(f"✅ Gemini supporte directement .{file_extension}")
                return {"extracted_text": response.text}
            except Exception as e:
                if "Unsupported MIME type" in str(e):
                    print(f"⚠️  Gemini ne supporte pas .{file_extension}, extraction de texte...")
                    # FALLBACK: Extraire le texte
                    text_content = extract_text_from_word(file_content, file_extension)
                    response = _client.models.generate_content(
                        model=GEMINI_MODEL,
                        contents=[prompt, text_content]
                    )
                    return {"extracted_text": response.text}
                else:
                    raise

        elif file_extension in ['pptx', 'ppt']:
            # Essayer d'abord avec le MIME type PowerPoint
            mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            try:
                response = _client.models.generate_content(
                    model=GEMINI_MODEL,
                    contents=[prompt, types.Part.from_bytes(data=file_content, mime_type=mime_type)]
                )
                print(f"✅ Gemini supporte directement .{file_extension}")
                return {"extracted_text": response.text}
            except Exception as e:
                if "Unsupported MIME type" in str(e):
                    print(f"⚠️  Gemini ne supporte pas .{file_extension}, extraction de texte...")
                    # FALLBACK: Extraire le texte
                    text_content = extract_text_from_ppt(file_content, file_extension)
                    response = _client.models.generate_content(
                        model=GEMINI_MODEL,
                        contents=[prompt, text_content]
                    )
                    return {"extracted_text": response.text}
                else:
                    raise
        else:
            raise ValueError(f"Type de fichier non supporté: {file_extension}")

    except Exception as e:
        raise ValueError(f"Erreur lors de l'extraction du fichier: {str(e)}")


def extract_text_from_word(file_content: bytes, extension: str) -> str:
    """Extraire le texte d'un fichier Word"""
    try:
        if extension == 'docx':
            from docx import Document
            doc = Document(io.BytesIO(file_content))
            text_content = '\n'.join([para.text for para in doc.paragraphs if para.text.strip()])
        else:  # .doc
            text_content = file_content.decode('utf-8', errors='ignore')

        if not text_content or len(text_content.strip()) < 50:
            raise ValueError("Le document Word semble vide ou illisible")

        return text_content
    except Exception as e:
        raise ValueError(f"Impossible d'extraire le texte du fichier Word: {str(e)}")


def extract_text_from_ppt(file_content: bytes, extension: str) -> str:
    """Extraire le texte d'un fichier PowerPoint"""
    try:
        if extension == 'pptx':
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_content))
            text_content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
            text_content = '\n'.join(text_content)
        else:  # .ppt
            text_content = file_content.decode('utf-8', errors='ignore')

        if not text_content or len(text_content.strip()) < 50:
            raise ValueError("La présentation PowerPoint semble vide ou illisible")

        return text_content
    except Exception as e:
        raise ValueError(f"Impossible d'extraire le texte du fichier PowerPoint: {str(e)}")


async def extract_from_text(text: str) -> dict:
    """
    Extraire les informations d'un texte brut

    Args:
        text: Texte brut contenant les informations

    Returns:
        dict: Contenu extrait
    """
    return {"extracted_text": text}


async def generate_format2_sfd(extracted_content: dict) -> Format2SFD:
    """
    Générer un SFD au format 2 (Agile) à partir du contenu extrait

    Args:
        extracted_content: Contenu extrait par Gemini

    Returns:
        Format2SFD: Objet SFD structuré
    """
    prompt = f"""
    À partir du texte suivant, génère une spécification fonctionnelle agile au format JSON strictement conforme à cette structure:

    {{
        "nom_projet": "string",
        "version": "string",
        "date": "string (format YYYY-MM-DD)",
        "product_owner": "string",
        "scrum_master": "string",
        "vision_produit": {{
            "probleme": "string",
            "utilisateurs_cibles": ["string"],
            "valeur_apportee": "string",
            "objectifs": ["string"]
        }},
        "epics": [
            {{
                "id": "string (ex: EP001)",
                "nom": "string",
                "description": "string",
                "objectif": "string",
                "user_stories": [
                    {{
                        "id": "string (ex: US001)",
                        "titre": "string",
                        "en_tant_que": "string",
                        "je_veux": "string",
                        "afin_de": "string",
                        "criteres_acceptation": ["string"],
                        "regles_metier": ["string"],
                        "priorite": "string (Haute/Moyenne/Basse)",
                        "estimation": "string",
                        "statut": "string"
                    }}
                ]
            }}
        ],
        "regles_metier": [
            {{
                "id": "string (ex: R001)",
                "nom": "string",
                "description": "string",
                "impact": ["string"]
            }}
        ],
        "modele_data": {{
            "entites": [{{"nom": "string", "attributs": "string"}}],
            "relations": ["string"]
        }},
        "workflows": [
            {{
                "nom": "string",
                "etapes": ["string"],
                "description": "string"
            }}
        ],
        "definition_of_done": ["string"],
        "notes": "string"
    }}

    Texte à analyser:
    {extracted_content.get('extracted_text', '')}

    IMPORTANT:
    - Retourne UNIQUEMENT le JSON, sans texte avant ou après
    - Respecte EXACTEMENT la structure ci-dessus
    - Utilise la date d'aujourd'hui: {datetime.now().strftime('%Y-%m-%d')}
    - Formule les user stories au format "En tant que... je veux... afin de..."
    - Si une information manque, utilise des valeurs par défaut raisonnables
    - Assure-toi que tous les IDs soient uniques et cohérents
    """

    response = _client.models.generate_content(
        model=GEMINI_MODEL,
        contents=prompt
    )

    # Nettoyer la réponse (enlever les balises markdown si présentes)
    json_text = response.text.strip()
    if json_text.startswith('```json'):
        json_text = json_text[7:]
    if json_text.startswith('```'):
        json_text = json_text[3:]
    if json_text.endswith('```'):
        json_text = json_text[:-3]
    json_text = json_text.strip()

    # Parser le JSON
    sfd_dict = json.loads(json_text)

    # Valider et créer l'objet Pydantic
    sfd = Format2SFD(**sfd_dict)

    return sfd


async def process_format2_extraction(
    file: Union[UploadFile, None] = None,
    text: Union[str, None] = None
) -> Format2SFD:
    """
    Processus complet d'extraction et génération pour Format 2

    Args:
        file: Fichier uploadé (optionnel)
        text: Texte brut (optionnel)

    Returns:
        Format2SFD: SFD structuré
    """
    if file:
        extracted_content = await extract_from_file(file)
    elif text:
        extracted_content = await extract_from_text(text)
    else:
        raise ValueError("Vous devez fournir soit un fichier, soit du texte")

    sfd = await generate_format2_sfd(extracted_content)

    return sfd
